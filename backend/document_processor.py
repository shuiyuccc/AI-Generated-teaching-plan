"""
文档处理模块（桩模块 - 用于开发调试）
提供参考文档的上传和解析功能
"""
import os
import logging

logger = logging.getLogger(__name__)

# 支持的文件扩展名
SUPPORTED_EXTENSIONS = {'.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.txt', '.pdf'}


def extract_document_content(file_path):
    """
    提取文档内容

    Args:
        file_path: 文档路径

    Returns:
        文档文本内容，失败返回 None
    """
    if not os.path.exists(file_path):
        logger.error(f"文件不存在: {file_path}")
        return None

    ext = os.path.splitext(file_path)[1].lower()

    if ext not in SUPPORTED_EXTENSIONS:
        logger.error(f"不支持的文件格式: {ext}")
        return None

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(file_path)
            return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])

        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
            return '\n'.join(texts)

        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join([str(c) for c in row if c is not None])
                    if row_text.strip():
                        texts.append(row_text)
            return '\n'.join(texts)

        elif ext == '.pdf':
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    texts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            texts.append(text)
                    return '\n'.join(texts)
            except ImportError:
                logger.warning("pdfplumber 未安装，PDF 解析不可用")
                return f"[PDF文件] {os.path.basename(file_path)}（需安装 pdfplumber）"

        return None

    except Exception as e:
        logger.error(f"文档解析失败: {e}")
        return None


def get_document_summary(file_path, max_chars=500):
    """
    获取文档摘要

    Args:
        file_path: 文档路径
        max_chars: 最大字符数

    Returns:
        文档摘要文本
    """
    content = extract_document_content(file_path)
    if content:
        return content[:max_chars] + ('...' if len(content) > max_chars else '')
    return None
